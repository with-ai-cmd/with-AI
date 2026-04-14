#!/usr/bin/env python3
"""AI SHINE 営業資料 2種類生成スクリプト（白ベース）"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ===== カラーパレット（白ベース） =====
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_WHITE = RGBColor(0xF8, 0xF9, 0xFA)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
GOLD = RGBColor(0xD4, 0xA5, 0x37)
BLUE = RGBColor(0x1E, 0x3A, 0x5F)
LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF5)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
RED = RGBColor(0xE0, 0x40, 0x40)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
CARD_BG = RGBColor(0xF5, 0xF7, 0xFA)
ACCENT_BLUE = RGBColor(0x29, 0x62, 0xFF)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                font_name="Meiryo"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, font_size=16,
                   color=DARK_TEXT, line_spacing=1.5, bold=False):
    """lines: list of (text, optional_overrides_dict)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if isinstance(line_data, str):
            txt, overrides = line_data, {}
        else:
            txt, overrides = line_data[0], line_data[1] if len(line_data) > 1 else {}
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(overrides.get("size", font_size))
        p.font.bold = overrides.get("bold", bold)
        p.font.color.rgb = overrides.get("color", color)
        p.font.name = overrides.get("font", "Meiryo")
        p.alignment = overrides.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(overrides.get("space_after", 4))
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color=CARD_BG,
                     line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_accent_line(slide, left, top, width, color=GOLD):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   left, top, width, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ====================================================================
# スライド作成関数
# ====================================================================

def slide_cover(prs, version):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, WHITE)

    # アクセントライン上
    add_accent_line(slide, Inches(3), Inches(1.8), Inches(7))

    # メインタイトル
    add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
                "会社の未来が輝く", font_size=44, bold=True, color=DARK_TEXT,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(3.0), Inches(11), Inches(1.5),
                "AI SHINE", font_size=72, bold=True, color=GOLD,
                alignment=PP_ALIGN.CENTER)

    # アクセントライン下
    add_accent_line(slide, Inches(3), Inches(4.5), Inches(7))

    # サブタイトル
    if version == "solo":
        subtitle = "〜 5年後の明暗を分ける、社長専属AIエージェント構築サービス 〜"
    else:
        subtitle = "〜 AIエージェントで実現する、次世代の業務体制 〜"

    add_textbox(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.8),
                subtitle, font_size=22, color=MEDIUM_GRAY,
                alignment=PP_ALIGN.CENTER)

    # 会社名
    add_textbox(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
                "Confidential | with-AI 株式会社", font_size=14, color=MEDIUM_GRAY,
                alignment=PP_ALIGN.CENTER)


def slide_3walls(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                "経営者が今、直面している 3つの壁", font_size=36, bold=True, color=DARK_TEXT)
    add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(5))

    cards = [
        ("01", "深刻な人手不足", "採用しても定着しない。\n育成コストだけが積み上がり、\n現場は常に人が足りない状態が続く。"),
        ("02", "社長の時間がない", "雑務・事務処理に追われ、\n本来注力すべき経営判断や\n新規事業に集中できない。"),
        ("03", "DXの遅れが致命傷に", "ツールは導入したが使いこなせず\n形骸化。競合はAIで武装し、\n差は日に日に開いている。"),
    ]
    for i, (num, title, desc) in enumerate(cards):
        x = Inches(0.8 + i * 4.1)
        add_rounded_rect(slide, x, Inches(1.8), Inches(3.7), Inches(4.2),
                         fill_color=LIGHT_BLUE, line_color=ACCENT_BLUE)
        # 番号バッジ
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3),
                                       Inches(2.1), Inches(0.5), Inches(0.5))
        badge.fill.solid()
        badge.fill.fore_color.rgb = ACCENT_BLUE
        badge.line.fill.background()
        badge.text_frame.paragraphs[0].text = num
        badge.text_frame.paragraphs[0].font.size = Pt(16)
        badge.text_frame.paragraphs[0].font.bold = True
        badge.text_frame.paragraphs[0].font.color.rgb = WHITE
        badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_textbox(slide, x + Inches(0.2), Inches(2.8), Inches(3.3), Inches(0.6),
                    title, font_size=22, bold=True, color=DARK_TEXT)
        add_textbox(slide, x + Inches(0.2), Inches(3.5), Inches(3.3), Inches(2.0),
                    desc, font_size=15, color=MEDIUM_GRAY)

    add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.6),
                "これらの課題を、テクノロジーの力で根本から解決する方法があります。",
                font_size=18, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)


def slide_ai_evolution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                "AIは「相談相手」から「実行する部下」へ進化した", font_size=34, bold=True, color=DARK_TEXT)
    add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(6))

    # Before card
    add_rounded_rect(slide, Inches(0.8), Inches(1.8), Inches(5.2), Inches(4.5),
                     fill_color=LIGHT_GRAY)
    add_textbox(slide, Inches(1.0), Inches(1.9), Inches(4.8), Inches(0.4),
                "これまでのAI", font_size=14, color=MEDIUM_GRAY)
    add_textbox(slide, Inches(1.0), Inches(2.3), Inches(4.8), Inches(0.6),
                '優秀な「相談相手」', font_size=26, bold=True, color=DARK_TEXT)
    before_items = [
        "● 質問すれば答えてくれる",
        "● しかし「実行」はしてくれない",
        "● コピー＆ペーストは人間の仕事",
        "● 1つの会話で1つの作業のみ",
        "",
        "→ 結局、人間の手間は減らない",
    ]
    add_multi_text(slide, Inches(1.0), Inches(3.2), Inches(4.8), Inches(3.0),
                   [(t, {"color": RED if "→" in t else MEDIUM_GRAY}) for t in before_items],
                   font_size=16)

    # Arrow
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   Inches(6.2), Inches(3.5), Inches(0.8), Inches(0.6))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GOLD
    arrow.line.fill.background()

    # After card
    add_rounded_rect(slide, Inches(7.2), Inches(1.8), Inches(5.5), Inches(4.5),
                     fill_color=LIGHT_BLUE, line_color=ACCENT_BLUE)
    add_textbox(slide, Inches(7.4), Inches(1.9), Inches(5.0), Inches(0.4),
                "今のAI ＝ AIエージェント", font_size=14, color=ACCENT_BLUE)
    add_textbox(slide, Inches(7.4), Inches(2.3), Inches(5.0), Inches(0.6),
                '実行する「優秀な部下」', font_size=26, bold=True, color=DARK_TEXT)
    after_items = [
        "● PCを操作し、業務を最後まで実行",
        "● 複数ツールを自律的に連携",
        "● 5人のチームを並列で同時稼働",
        "● 24時間365日、文句なく働き続ける",
        "",
        "→ 人間は指示を出すだけでOK",
    ]
    add_multi_text(slide, Inches(7.4), Inches(3.2), Inches(5.0), Inches(3.0),
                   [(t, {"color": GREEN if "→" in t else DARK_TEXT}) for t in after_items],
                   font_size=16)

    # Bottom message
    add_rounded_rect(slide, Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.7),
                     fill_color=GOLD)
    add_textbox(slide, Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.7),
                "この「AIエージェント」を自社に組み込めるかどうかが、5年後の経営格差を決定づけます。",
                font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)


def slide_5year_comparison(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                "5年後、あなたの会社はどちらですか？", font_size=36, bold=True, color=DARK_TEXT)
    add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(5))

    # AI導入企業
    add_rounded_rect(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5),
                     fill_color=LIGHT_BLUE, line_color=ACCENT_BLUE)
    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(0.6),
                "AI導入企業", font_size=26, bold=True, color=ACCENT_BLUE,
                alignment=PP_ALIGN.CENTER)
    ai_items = [
        "◆  少人数で高収益体制を実現",
        "◆  社長は意思決定に 100%集中",
        "◆  ルーティン業務は AIが自動処理",
        "◆  採用難の影響を受けない",
        "◆  競合に対して圧倒的な優位性",
    ]
    add_multi_text(slide, Inches(1.2), Inches(2.8), Inches(4.8), Inches(3.5),
                   ai_items, font_size=18, bold=True, color=DARK_TEXT)

    # 未導入企業
    add_rounded_rect(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.5),
                     fill_color=LIGHT_GRAY)
    add_textbox(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(0.6),
                "未導入企業", font_size=26, bold=True, color=MEDIUM_GRAY,
                alignment=PP_ALIGN.CENTER)
    no_items = [
        "▪  人件費の増大が利益を圧迫",
        "▪  属人化が進みリスクが増大",
        "▪  社長が雑務から解放されない",
        "▪  AI活用企業との格差が拡大",
        "▪  市場での競争力が低下",
    ]
    add_multi_text(slide, Inches(7.4), Inches(2.8), Inches(4.8), Inches(3.5),
                   no_items, font_size=18, color=MEDIUM_GRAY)

    add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.6),
                '手遅れになる前の "今"、最初の一歩を踏み出しませんか？',
                font_size=20, bold=True, color=RED, alignment=PP_ALIGN.CENTER)


def slide_what_is(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_accent_line(slide, Inches(3.5), Inches(2.0), Inches(6.3))

    add_textbox(slide, Inches(1), Inches(2.3), Inches(11), Inches(1.0),
                "AI SHINE とは？", font_size=48, bold=True, color=DARK_TEXT,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(1.0),
                "あなた専属のAI社員チームを構築し、\n育て続ける伴走型サービス",
                font_size=24, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

    labels = ["構 築", "教 育", "伴 走"]
    for i, label in enumerate(labels):
        x = Inches(3.2 + i * 2.5)
        add_rounded_rect(slide, x, Inches(5.0), Inches(2.0), Inches(0.8),
                         fill_color=ACCENT_BLUE, line_color=GOLD)
        add_textbox(slide, x, Inches(5.05), Inches(2.0), Inches(0.7),
                    label, font_size=22, bold=True, color=WHITE,
                    alignment=PP_ALIGN.CENTER)

    add_accent_line(slide, Inches(3.5), Inches(6.0), Inches(6.3))


def slide_3changes(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                "AI SHINE導入で起こる 3つの劇的変化", font_size=34, bold=True, color=DARK_TEXT)
    add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(5.5))

    changes = [
        ("01", "24時間365日、稼働し続ける環境",
         "人間の社員には労働時間の規制がありますが、AI社員に制限はありません。\nあなたが寝ている夜間も休日も、AIは業務をこなし続けます。\n朝起きれば完璧なレポートやスケジュールが完成しています。"),
        ("02", "「5人のAIチーム」で人手不足を完全解消",
         "1回の指示で、リサーチ担当・戦略担当・資料作成担当など\n5人のAIチームを並列で同時に動かせます。\n採用難の時代に、圧倒的な処理能力を即座に配置できます。"),
        ("03", "利益率を最大化し、企業の未来を輝かせる",
         "ルーティン業務をすべてAI社員に任せることで、\n社長や人間の社員は意思決定やクリエイティブな業務に100%集中。\nコストを抑えながら利益率を最大化します。"),
    ]

    for i, (num, title, desc) in enumerate(changes):
        y = Inches(1.6 + i * 1.9)
        add_rounded_rect(slide, Inches(0.8), y, Inches(11.7), Inches(1.7),
                         fill_color=LIGHT_BLUE if i % 2 == 0 else CARD_BG,
                         line_color=ACCENT_BLUE if i % 2 == 0 else None)
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       Inches(1.1), y + Inches(0.2),
                                       Inches(0.5), Inches(0.5))
        badge.fill.solid()
        badge.fill.fore_color.rgb = ACCENT_BLUE
        badge.line.fill.background()
        badge.text_frame.paragraphs[0].text = num
        badge.text_frame.paragraphs[0].font.size = Pt(16)
        badge.text_frame.paragraphs[0].font.bold = True
        badge.text_frame.paragraphs[0].font.color.rgb = WHITE
        badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_textbox(slide, Inches(2.0), y + Inches(0.15), Inches(10.0), Inches(0.5),
                    title, font_size=22, bold=True, color=GOLD)
        add_textbox(slide, Inches(2.0), y + Inches(0.65), Inches(10.0), Inches(1.0),
                    desc, font_size=14, color=MEDIUM_GRAY)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                "AI SHINEの全体像", font_size=36, bold=True, color=DARK_TEXT)
    add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(3.5))

    # 社長
    add_rounded_rect(slide, Inches(0.5), Inches(2.8), Inches(2.0), Inches(1.8),
                     fill_color=CARD_BG, line_color=GOLD)
    add_textbox(slide, Inches(0.5), Inches(3.0), Inches(2.0), Inches(0.5),
                "👤 社長（あなた）", font_size=14, bold=True, color=DARK_TEXT,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.5), Inches(2.0), Inches(0.5),
                "指示を出すだけ", font_size=13, color=MEDIUM_GRAY,
                alignment=PP_ALIGN.CENTER)

    # Arrow
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                    Inches(2.7), Inches(3.4), Inches(0.7), Inches(0.5))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = GOLD
    arrow1.line.fill.background()

    # 親エージェント
    add_rounded_rect(slide, Inches(3.6), Inches(2.5), Inches(2.5), Inches(2.4),
                     fill_color=LIGHT_BLUE, line_color=ACCENT_BLUE)
    add_textbox(slide, Inches(3.6), Inches(2.7), Inches(2.5), Inches(0.5),
                "🧠 親エージェント", font_size=15, bold=True, color=ACCENT_BLUE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(3.6), Inches(3.3), Inches(2.5), Inches(1.0),
                "指示を分解し\nサブエージェントに\n自律的に配分",
                font_size=13, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

    # Arrow
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                    Inches(6.3), Inches(3.4), Inches(0.7), Inches(0.5))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = ACCENT_BLUE
    arrow2.line.fill.background()

    # Sub agents
    agents = [
        ("📋 秘書", "Google Drive"),
        ("💰 経理", "freee"),
        ("🔍 リサーチ", "WordPress"),
        ("✍️ ライター", "Slack"),
        ("📊 監査", "Web検索"),
    ]
    for i, (agent, tool) in enumerate(agents):
        y = Inches(1.2 + i * 1.2)
        add_rounded_rect(slide, Inches(7.3), y, Inches(2.2), Inches(0.9),
                         fill_color=CARD_BG, line_color=ACCENT_BLUE)
        add_textbox(slide, Inches(7.3), y + Inches(0.15), Inches(2.2), Inches(0.6),
                    agent, font_size=15, bold=True, color=DARK_TEXT,
                    alignment=PP_ALIGN.CENTER)

        # small arrow
        arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                     Inches(9.7), y + Inches(0.25), Inches(0.5), Inches(0.35))
        arr.fill.solid()
        arr.fill.fore_color.rgb = MEDIUM_GRAY
        arr.line.fill.background()

        add_rounded_rect(slide, Inches(10.4), y, Inches(2.2), Inches(0.9),
                         fill_color=LIGHT_GRAY)
        add_textbox(slide, Inches(10.4), y + Inches(0.15), Inches(2.2), Inches(0.6),
                    tool, font_size=14, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)


def slide_case(prs, case_num, title, problem, flow_items=None, highlight_text=None, before_after=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.2), Inches(3), Inches(0.4),
                f"活用事例 ④" if case_num == 4 else f"活用事例 {['①','②','③','④'][case_num-1]}",
                font_size=14, color=ACCENT_BLUE)
    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                title, font_size=30, bold=True, color=DARK_TEXT)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(5))

    # Problem bar
    add_rounded_rect(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.8),
                     fill_color=LIGHT_GRAY)
    add_textbox(slide, Inches(1.0), Inches(1.55), Inches(11.3), Inches(0.7),
                f"😩 悩み：{problem}", font_size=15, color=MEDIUM_GRAY)

    if flow_items:
        for i, item in enumerate(flow_items):
            x = Inches(0.5 + i * 2.6)
            add_rounded_rect(slide, x, Inches(2.8), Inches(2.2), Inches(1.5),
                             fill_color=LIGHT_BLUE, line_color=ACCENT_BLUE)
            add_textbox(slide, x, Inches(3.0), Inches(2.2), Inches(1.2),
                        item, font_size=14, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
            if i < len(flow_items) - 1:
                arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                             x + Inches(2.3), Inches(3.3),
                                             Inches(0.25), Inches(0.3))
                arr.fill.solid()
                arr.fill.fore_color.rgb = ACCENT_BLUE
                arr.line.fill.background()

    if highlight_text:
        add_rounded_rect(slide, Inches(0.8), Inches(4.8), Inches(11.7), Inches(1.8),
                         fill_color=CARD_BG, line_color=GOLD)
        add_textbox(slide, Inches(1.2), Inches(4.9), Inches(11.0), Inches(1.6),
                    highlight_text, font_size=15, color=DARK_TEXT)

    if before_after:
        bf, af = before_after
        add_rounded_rect(slide, Inches(0.8), Inches(2.8), Inches(5.2), Inches(3.5),
                         fill_color=LIGHT_GRAY)
        add_textbox(slide, Inches(0.8), Inches(2.9), Inches(5.2), Inches(0.5),
                    "BEFORE（従来の方法）", font_size=16, color=MEDIUM_GRAY,
                    alignment=PP_ALIGN.CENTER)
        add_multi_text(slide, Inches(1.2), Inches(3.5), Inches(4.5), Inches(2.5),
                       bf, font_size=16, color=MEDIUM_GRAY)

        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       Inches(6.2), Inches(4.2), Inches(0.7), Inches(0.5))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GOLD
        arrow.line.fill.background()

        add_rounded_rect(slide, Inches(7.2), Inches(2.8), Inches(5.3), Inches(3.5),
                         fill_color=LIGHT_BLUE, line_color=ACCENT_BLUE)
        add_textbox(slide, Inches(7.2), Inches(2.9), Inches(5.3), Inches(0.5),
                    "AFTER（AI SHINE）", font_size=16, color=ACCENT_BLUE,
                    alignment=PP_ALIGN.CENTER)
        add_multi_text(slide, Inches(7.6), Inches(3.5), Inches(4.5), Inches(2.5),
                       af, font_size=16, color=DARK_TEXT, bold=True)


def slide_withai_orgchart(prs):
    """with-AI社の実際のAI組織図 — 具体的なイメージを伝える"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.4),
                "実例：with-AI社の場合", font_size=14, color=ACCENT_BLUE)
    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                "私たちは実際にこの体制で会社を回しています", font_size=34, bold=True, color=DARK_TEXT)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(6))

    # ===== 社長（勝又） =====
    add_rounded_rect(slide, Inches(5.2), Inches(1.5), Inches(2.8), Inches(0.9),
                     fill_color=GOLD)
    add_textbox(slide, Inches(5.2), Inches(1.55), Inches(2.8), Inches(0.8),
                "👤 代表取締役（人間）\n指示を出すだけ", font_size=14, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    # 縦線 社長→CAIO
    vline1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(6.55), Inches(2.4), Pt(3), Inches(0.4))
    vline1.fill.solid()
    vline1.fill.fore_color.rgb = MEDIUM_GRAY
    vline1.line.fill.background()

    # ===== CAIO =====
    add_rounded_rect(slide, Inches(4.7), Inches(2.8), Inches(3.8), Inches(0.9),
                     fill_color=DARK_TEXT)
    add_textbox(slide, Inches(4.7), Inches(2.85), Inches(3.8), Inches(0.8),
                "🧠 CAIO（最高AI責任者）\nAI戦略の司令塔・全体統括", font_size=14, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    # 縦線 CAIO→秘書
    vline2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(6.55), Inches(3.7), Pt(3), Inches(0.4))
    vline2.fill.solid()
    vline2.fill.fore_color.rgb = MEDIUM_GRAY
    vline2.line.fill.background()

    # ===== 秘書AI =====
    add_rounded_rect(slide, Inches(4.7), Inches(4.1), Inches(3.8), Inches(0.9),
                     fill_color=ACCENT_BLUE)
    add_textbox(slide, Inches(4.7), Inches(4.15), Inches(3.8), Inches(0.8),
                "📋 秘書AI（親エージェント）\n指示を分解し各部門に配分", font_size=14, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    # 横線 秘書→各部門
    hline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(1.6), Inches(5.25), Inches(10.0), Pt(3))
    hline.fill.solid()
    hline.fill.fore_color.rgb = MEDIUM_GRAY
    hline.line.fill.background()

    # 縦線 秘書→横線
    vline3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(6.55), Inches(5.0), Pt(3), Inches(0.25))
    vline3.fill.solid()
    vline3.fill.fore_color.rgb = MEDIUM_GRAY
    vline3.line.fill.background()

    # ===== 4つの部門 =====
    departments = [
        ("CTO", "最高技術責任者", "技術・開発\nシステム構築",
         ["GitHub", "Claude Code", "API連携"]),
        ("CPO", "最高プロダクト責任者", "サービス構築\n資料作成",
         ["Google Slides", "Notion", "Drive"]),
        ("CFO", "最高財務責任者", "経理・財務\n請求・税務",
         ["freee", "GMO銀行", "Google Drive"]),
        ("CLO", "最高法務責任者", "契約書・規約\n法令リサーチ",
         ["契約書DB", "法令検索", "PDF生成"]),
    ]

    for i, (title, subtitle, desc, tools) in enumerate(departments):
        x = Inches(0.3 + i * 3.2)

        # 縦線 横線→部門
        vl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    x + Inches(1.3), Inches(5.25), Pt(3), Inches(0.25))
        vl.fill.solid()
        vl.fill.fore_color.rgb = MEDIUM_GRAY
        vl.line.fill.background()

        # 部門カード
        add_rounded_rect(slide, x, Inches(5.5), Inches(2.8), Inches(1.7),
                         fill_color=LIGHT_BLUE, line_color=ACCENT_BLUE)
        add_textbox(slide, x, Inches(5.55), Inches(2.8), Inches(0.4),
                    f"🏢 {title}", font_size=20, bold=True, color=ACCENT_BLUE,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(5.9), Inches(2.8), Inches(0.3),
                    subtitle, font_size=11, color=MEDIUM_GRAY,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(6.2), Inches(2.8), Inches(0.5),
                    desc, font_size=12, color=DARK_TEXT,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(6.7), Inches(2.8), Inches(0.4),
                    " | ".join(tools), font_size=9, color=MEDIUM_GRAY,
                    alignment=PP_ALIGN.CENTER)


def slide_skill_system(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                "「いつものあれ、やって」で動く魔法の仕組み", font_size=34, bold=True, color=DARK_TEXT)
    add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(5.5))

    steps = [
        ("STEP 1", "一度、業務を\nAIに教える", "初回だけ丁寧に\n手順を指示"),
        ("STEP 2", "「スキル」として\n自動保存", "AIが手順を\n記憶・最適化"),
        ("STEP 3", "次回からは\nワンコマンド", "「いつものあれ」\nで即実行"),
    ]
    for i, (step, title, desc) in enumerate(steps):
        x = Inches(1.0 + i * 4.0)
        add_rounded_rect(slide, x, Inches(1.8), Inches(3.5), Inches(3.2),
                         fill_color=LIGHT_BLUE if i == 2 else CARD_BG,
                         line_color=ACCENT_BLUE if i == 2 else None)
        add_textbox(slide, x, Inches(1.9), Inches(3.5), Inches(0.4),
                    step, font_size=14, bold=True, color=ACCENT_BLUE,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(2.4), Inches(3.5), Inches(1.2),
                    title, font_size=22, bold=True, color=DARK_TEXT,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(3.6), Inches(3.5), Inches(1.0),
                    desc, font_size=15, color=MEDIUM_GRAY,
                    alignment=PP_ALIGN.CENTER)

        if i < 2:
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                         x + Inches(3.6), Inches(3.2),
                                         Inches(0.35), Inches(0.4))
            arr.fill.solid()
            arr.fill.fore_color.rgb = GOLD
            arr.line.fill.background()

    # Command examples
    cmds = ['"見積もり作って"', '"いつものリサーチして"', '"競合チェックして"', '"ブログ書いて"']
    for i, cmd in enumerate(cmds):
        x = Inches(1.0 + i * 3.1)
        add_rounded_rect(slide, x, Inches(5.3), Inches(2.8), Inches(0.6),
                         fill_color=ACCENT_BLUE)
        add_textbox(slide, x, Inches(5.3), Inches(2.8), Inches(0.6),
                    cmd, font_size=15, bold=True, color=WHITE,
                    alignment=PP_ALIGN.CENTER)

    add_rounded_rect(slide, Inches(1.5), Inches(6.2), Inches(10.3), Inches(0.8),
                     fill_color=CARD_BG, line_color=GOLD)
    add_textbox(slide, Inches(1.5), Inches(6.25), Inches(10.3), Inches(0.7),
                "一人社長であっても、文句を言わず、ミスをせず、24時間働き続ける\n「最強のチーム」を PCの中に持つことができるのです。",
                font_size=16, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)


def slide_comparison_table(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                "なぜAI SHINEなのか？", font_size=36, bold=True, color=DARK_TEXT)
    add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(4))

    headers = ["", "ChatGPT等を\n自分で使う", "SaaSツール\n導入", "AI SHINE"]
    rows = [
        ("カスタマイズ性", "× 汎用的", "△ 設定範囲内", "◎ 御社専用に構築"),
        ("業務の自動実行", "× 手動コピペ", "△ 単一業務のみ", "◎ 複数業務を並列自動化"),
        ("継続的な進化", "× 自力で学習", "× アプデ待ち", "◎ 毎月の伴走で成長"),
        ("導入ハードル", "△ 知識が必要", "○ 比較的低い", "◎ すべて代行"),
        ("コスト効率", "○ 低コスト", "△ 月額課金", "◎ ROI最大化"),
    ]

    # Header row
    for j, h in enumerate(headers):
        x = Inches(0.5 + j * 3.1)
        color = ACCENT_BLUE if j == 3 else LIGHT_GRAY
        add_rounded_rect(slide, x, Inches(1.5), Inches(2.9), Inches(0.8), fill_color=color)
        add_textbox(slide, x, Inches(1.5), Inches(2.9), Inches(0.8),
                    h, font_size=15, bold=True,
                    color=WHITE if j == 3 else DARK_TEXT,
                    alignment=PP_ALIGN.CENTER)

    for i, (label, c1, c2, c3) in enumerate(rows):
        y = Inches(2.5 + i * 0.85)
        vals = [label, c1, c2, c3]
        for j, v in enumerate(vals):
            x = Inches(0.5 + j * 3.1)
            bg = LIGHT_BLUE if j == 3 else (CARD_BG if i % 2 == 0 else WHITE)
            add_rounded_rect(slide, x, y, Inches(2.9), Inches(0.7), fill_color=bg)
            fc = ACCENT_BLUE if j == 3 else (DARK_TEXT if j == 0 else MEDIUM_GRAY)
            add_textbox(slide, x, y + Inches(0.1), Inches(2.9), Inches(0.5),
                        v, font_size=14, bold=(j == 0 or j == 3), color=fc,
                        alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5),
                "AI SHINEは「ツール導入」ではなく「AI社員の育成・伴走」。だから、成果が出続けます。",
                font_size=17, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)


def slide_testimonials(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                "導入企業様の声", font_size=36, bold=True, color=DARK_TEXT)
    add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(3.5))

    testimonials = [
        ("製造業 A社 代表取締役",
         "「月末の経理処理が丸2日かかっていたのが、AIのおかげで確認ボタンを押すだけに。\n空いた時間で新規営業に集中でき、売上が前年比 120%に伸びました。」"),
        ("コンサル業 B社 代表",
         "「提案書の作成スピードが劇的に変わりました。商談翌日には高品質な提案書が\n完成しているので、受注率が明らかに上がっています。」"),
        ("不動産業 C社 社長",
         "「競合分析に外注で30万円かけていたのが、今は URLを入れるだけ。\nしかもスピードも精度も外注以上。もっと早く導入すべきでした。」"),
    ]

    for i, (company, quote) in enumerate(testimonials):
        y = Inches(1.6 + i * 1.7)
        add_rounded_rect(slide, Inches(0.8), y, Inches(11.7), Inches(1.5),
                         fill_color=CARD_BG, line_color=ACCENT_BLUE if i == 0 else None)
        add_textbox(slide, Inches(1.2), y + Inches(0.15), Inches(10), Inches(0.4),
                    company, font_size=16, bold=True, color=GOLD)
        add_textbox(slide, Inches(1.2), y + Inches(0.55), Inches(10.8), Inches(0.9),
                    quote, font_size=14, color=MEDIUM_GRAY)

    add_rounded_rect(slide, Inches(2.5), Inches(6.5), Inches(8.3), Inches(0.6),
                     fill_color=LIGHT_BLUE, line_color=GOLD)
    add_textbox(slide, Inches(2.5), Inches(6.5), Inches(8.3), Inches(0.6),
                "※ 現在、先行導入モニター企業様を募集しております",
                font_size=16, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)


def slide_step1_build(prs, version):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(1), Inches(0.3), Inches(11), Inches(0.5),
                "STEP 1", font_size=18, bold=True, color=ACCENT_BLUE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(0.7), Inches(11), Inches(1.0),
                'まずはAI社員を「入社」させる', font_size=40, bold=True, color=DARK_TEXT,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.5),
                "— 初期構築プラン —", font_size=20, color=MEDIUM_GRAY,
                alignment=PP_ALIGN.CENTER)
    add_accent_line(slide, Inches(4), Inches(2.1), Inches(5.3))

    add_textbox(slide, Inches(1), Inches(2.3), Inches(11), Inches(0.5),
                "構築内容", font_size=22, bold=True, color=DARK_TEXT,
                alignment=PP_ALIGN.CENTER)

    items = [
        ("01", "AI環境の導入", "お客様のPCにAIエージェント環境を完全セットアップ"),
        ("02", "会社情報の構造化", "フォルダ構造化によるAIへのコンテキスト付与（会社のDNA注入）"),
        ("03", "基本スキルの実装", "秘書・経理・リサーチ等の基本業務スキルをAIに教育"),
        ("04", "外部ツール連携", "Google Drive / freee / Slack / WordPress等とのMCP連携構築"),
    ]
    for i, (num, title, desc) in enumerate(items):
        row = i // 2
        col = i % 2
        x = Inches(0.8 + col * 6.2)
        y = Inches(3.0 + row * 1.7)
        add_rounded_rect(slide, x, y, Inches(5.8), Inches(1.4),
                         fill_color=LIGHT_BLUE, line_color=ACCENT_BLUE)
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2),
                                       y + Inches(0.2), Inches(0.45), Inches(0.45))
        badge.fill.solid()
        badge.fill.fore_color.rgb = ACCENT_BLUE
        badge.line.fill.background()
        badge.text_frame.paragraphs[0].text = num
        badge.text_frame.paragraphs[0].font.size = Pt(14)
        badge.text_frame.paragraphs[0].font.bold = True
        badge.text_frame.paragraphs[0].font.color.rgb = WHITE
        badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_textbox(slide, x + Inches(0.8), y + Inches(0.15), Inches(4.8), Inches(0.5),
                    title, font_size=18, bold=True, color=GOLD)
        add_textbox(slide, x + Inches(0.8), y + Inches(0.6), Inches(4.8), Inches(0.6),
                    desc, font_size=13, color=MEDIUM_GRAY)

    if version == "solo":
        price_text = "通常価格  500,000 円（税別）"
    else:
        price_text = "通常価格  500,000 円（税別）"

    add_textbox(slide, Inches(1), Inches(6.7), Inches(11), Inches(0.6),
                price_text, font_size=22, bold=True, color=DARK_TEXT,
                alignment=PP_ALIGN.CENTER)


def slide_campaign(prs, version):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_rounded_rect(slide, Inches(1.5), Inches(0.8), Inches(10.3), Inches(6.0),
                     fill_color=CARD_BG, line_color=GOLD)

    if version == "solo":
        add_textbox(slide, Inches(1.5), Inches(1.0), Inches(10.3), Inches(0.6),
                    "🎯 期間限定キャンペーン", font_size=28, bold=True, color=DARK_TEXT,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10.3), Inches(0.5),
                    "初期構築プラン", font_size=20, color=MEDIUM_GRAY,
                    alignment=PP_ALIGN.CENTER)
        add_accent_line(slide, Inches(4), Inches(2.4), Inches(5.3))

        # Strikethrough price (simulated)
        add_textbox(slide, Inches(1.5), Inches(2.7), Inches(10.3), Inches(0.6),
                    "通常価格  500,000円", font_size=22, color=MEDIUM_GRAY,
                    alignment=PP_ALIGN.CENTER)

        # Discount badge
        badge = add_rounded_rect(slide, Inches(9.5), Inches(3.3), Inches(1.5), Inches(1.0),
                                 fill_color=RED)
        add_textbox(slide, Inches(9.5), Inches(3.35), Inches(1.5), Inches(0.9),
                    "40%\nOFF", font_size=20, bold=True, color=WHITE,
                    alignment=PP_ALIGN.CENTER)

        add_textbox(slide, Inches(1.5), Inches(3.5), Inches(8), Inches(0.4),
                    "特別価格", font_size=16, color=RED, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.5), Inches(3.9), Inches(8), Inches(1.2),
                    "300,000 円（税別）", font_size=56, bold=True, color=GOLD,
                    alignment=PP_ALIGN.CENTER)

        add_textbox(slide, Inches(1.5), Inches(5.3), Inches(10.3), Inches(0.4),
                    "※ 先着限定 ／ 本キャンペーンは予告なく終了する場合がございます",
                    font_size=13, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.5),
                    "まずはAIエージェントがどのようなものかを体感してください。",
                    font_size=17, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

    else:  # corporate
        add_textbox(slide, Inches(1.5), Inches(1.0), Inches(10.3), Inches(0.6),
                    "特別ご案内", font_size=28, bold=True, color=DARK_TEXT,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10.3), Inches(0.5),
                    "初期構築プラン", font_size=20, color=MEDIUM_GRAY,
                    alignment=PP_ALIGN.CENTER)
        add_accent_line(slide, Inches(4), Inches(2.4), Inches(5.3))

        add_textbox(slide, Inches(1.5), Inches(3.0), Inches(10.3), Inches(0.4),
                    "特別価格", font_size=16, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.5), Inches(3.4), Inches(10.3), Inches(1.2),
                    "150,000 円（税別）", font_size=56, bold=True, color=GOLD,
                    alignment=PP_ALIGN.CENTER)

        add_textbox(slide, Inches(1.5), Inches(4.8), Inches(10.3), Inches(0.8),
                    "※ ご紹介・パートナー企業様向けの特別価格です\n通常価格 500,000円のところ、70%OFFでご提供いたします",
                    font_size=15, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.5),
                    "まずはAIエージェントがどのようなものかを体感してください。",
                    font_size=17, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)


def slide_step2_intro(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.5),
                "STEP 2", font_size=18, bold=True, color=ACCENT_BLUE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(1.0), Inches(11), Inches(1.0),
                "AI社員を「育てる」— 月額伴走プラン", font_size=38, bold=True, color=DARK_TEXT,
                alignment=PP_ALIGN.CENTER)
    add_accent_line(slide, Inches(3.5), Inches(2.0), Inches(6.3))

    add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.6),
                "構築はゴールではなく、スタートです。", font_size=26, bold=True, color=GOLD,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.5), Inches(3.3), Inches(10), Inches(2.0),
                "AIエージェントの真価は「育て続ける」ことで発揮されます。\n\nAI技術は日々進化しています。\nどのような新しいAIが登場しても、御社の業務に最適な形で\nアップデートし続ける体制をご提供します。",
                font_size=18, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

    # Step 1 → Step 2
    add_rounded_rect(slide, Inches(2.5), Inches(5.3), Inches(3.5), Inches(1.5),
                     fill_color=CARD_BG)
    add_textbox(slide, Inches(2.5), Inches(5.35), Inches(3.5), Inches(0.4),
                "STEP 1 で構築", font_size=14, bold=True, color=ACCENT_BLUE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2.5), Inches(5.8), Inches(3.5), Inches(0.8),
                "AIエージェントが\nどのようなものかを知る", font_size=16, color=DARK_TEXT,
                alignment=PP_ALIGN.CENTER)

    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 Inches(6.3), Inches(5.8), Inches(0.7), Inches(0.5))
    arr.fill.solid()
    arr.fill.fore_color.rgb = GOLD
    arr.line.fill.background()

    add_rounded_rect(slide, Inches(7.3), Inches(5.3), Inches(3.5), Inches(1.5),
                     fill_color=LIGHT_BLUE, line_color=ACCENT_BLUE)
    add_textbox(slide, Inches(7.3), Inches(5.35), Inches(3.5), Inches(0.4),
                "STEP 2 で伴走", font_size=14, bold=True, color=ACCENT_BLUE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(7.3), Inches(5.8), Inches(3.5), Inches(0.8),
                "本当に自分の業務を\n自動化していく", font_size=16, bold=True, color=DARK_TEXT,
                alignment=PP_ALIGN.CENTER)


def slide_monthly_plans_solo(prs):
    """個人事業主向け：現行の3プラン（9,800 / 29,800 / 49,800）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(1), Inches(0.3), Inches(11), Inches(0.8),
                "月額伴走プラン", font_size=38, bold=True, color=DARK_TEXT,
                alignment=PP_ALIGN.CENTER)
    add_accent_line(slide, Inches(4), Inches(1.0), Inches(5.3))

    plans = [
        {
            "name": "LIGHT", "price": "¥9,800", "recommended": False,
            "features": [
                ("✅", "AI最新NEWS共有"),
                ("✅", "チャット相談（随時）"),
                ("—", "月1回オンライン相談"),
                ("—", "エージェント設計の提案"),
                ("—", "詳細な構築アドバイス"),
                ("—", "業務自動化の具体設計"),
            ],
            "tagline": "AIの最新動向を\nキャッチアップしたい方",
        },
        {
            "name": "STANDARD", "price": "¥29,800", "recommended": True,
            "features": [
                ("✅", "AI最新NEWS共有"),
                ("✅", "チャット相談（随時）"),
                ("✅", "月1回オンライン相談"),
                ("✅", "エージェント設計の提案"),
                ("—", "詳細な構築アドバイス"),
                ("—", "業務自動化の具体設計"),
            ],
            "tagline": "自社に合うエージェントを\n一緒に考えたい方",
        },
        {
            "name": "PREMIUM", "price": "¥49,800", "recommended": False,
            "features": [
                ("✅", "AI最新NEWS共有"),
                ("✅", "チャット相談（随時）"),
                ("✅", "月1回オンライン相談"),
                ("✅", "エージェント設計の提案"),
                ("✅", "詳細な構築アドバイス"),
                ("✅", "業務自動化の具体設計（簡易設計）"),
            ],
            "tagline": "本格的に業務自動化を\n推進したい方",
        },
    ]

    for i, plan in enumerate(plans):
        x = Inches(0.6 + i * 4.2)
        border = GOLD if plan["recommended"] else None
        bg = LIGHT_BLUE if plan["recommended"] else CARD_BG
        add_rounded_rect(slide, x, Inches(1.4), Inches(3.8), Inches(5.5),
                         fill_color=bg, line_color=border)

        if plan["recommended"]:
            rec_badge = add_rounded_rect(slide, x + Inches(1.2), Inches(1.2),
                                         Inches(1.4), Inches(0.35), fill_color=GOLD)
            add_textbox(slide, x + Inches(1.2), Inches(1.2), Inches(1.4), Inches(0.35),
                        "おすすめ", font_size=12, bold=True, color=WHITE,
                        alignment=PP_ALIGN.CENTER)

        add_textbox(slide, x, Inches(1.6), Inches(3.8), Inches(0.5),
                    plan["name"], font_size=24, bold=True, color=DARK_TEXT,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(2.0), Inches(3.8), Inches(0.3),
                    "月額（税別）", font_size=12, color=MEDIUM_GRAY,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(2.3), Inches(3.8), Inches(0.7),
                    plan["price"], font_size=38, bold=True, color=GOLD,
                    alignment=PP_ALIGN.CENTER)

        for j, (icon, feat) in enumerate(plan["features"]):
            y = Inches(3.2 + j * 0.38)
            c = DARK_TEXT if icon == "✅" else MEDIUM_GRAY
            add_textbox(slide, x + Inches(0.3), y, Inches(3.2), Inches(0.35),
                        f"{icon}  {feat}", font_size=13, color=c)

        # Tagline
        add_rounded_rect(slide, x + Inches(0.3), Inches(5.7), Inches(3.2), Inches(0.9),
                         fill_color=ACCENT_BLUE if plan["recommended"] else LIGHT_GRAY)
        add_textbox(slide, x + Inches(0.3), Inches(5.75), Inches(3.2), Inches(0.8),
                    plan["tagline"], font_size=13, bold=True,
                    color=WHITE if plan["recommended"] else DARK_TEXT,
                    alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.4),
                "※ すべてのプランは月払い・いつでもプラン変更可能です",
                font_size=13, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


def slide_monthly_plans_corporate(prs):
    """企業向け：3プラン（29,800 / 49,800 / 100,000）+ 追加人数"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(1), Inches(0.3), Inches(11), Inches(0.8),
                "月額伴走プラン", font_size=38, bold=True, color=DARK_TEXT,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(0.9), Inches(11), Inches(0.4),
                "基本：1〜5名", font_size=16, color=MEDIUM_GRAY,
                alignment=PP_ALIGN.CENTER)
    add_accent_line(slide, Inches(4), Inches(1.4), Inches(5.3))

    plans = [
        {
            "name": "LIGHT", "price": "¥29,800", "recommended": False,
            "subtitle": "相談プラン",
            "features": [
                ("✅", "AI最新NEWS共有"),
                ("✅", "チャット相談（随時）"),
                ("✅", "月1回オンライン相談"),
                ("—", "エージェント設計の提案"),
                ("—", "構築サポート"),
            ],
            "tagline": "まずはAI活用の\n相談から始めたい",
        },
        {
            "name": "STANDARD", "price": "¥49,800", "recommended": True,
            "subtitle": "提案プラン",
            "features": [
                ("✅", "AI最新NEWS共有"),
                ("✅", "チャット相談（随時）"),
                ("✅", "月1回オンライン相談"),
                ("✅", "エージェント設計の提案"),
                ("—", "構築サポート"),
            ],
            "tagline": "自社に最適な\nAI活用を設計したい",
        },
        {
            "name": "PREMIUM", "price": "¥100,000", "recommended": False,
            "subtitle": "構築サポートプラン",
            "features": [
                ("✅", "AI最新NEWS共有"),
                ("✅", "チャット相談（随時）"),
                ("✅", "月1回オンライン相談"),
                ("✅", "エージェント設計の提案"),
                ("✅", "エージェント構築サポート"),
            ],
            "tagline": "本格的にAI業務自動化を\n実装・推進したい",
        },
    ]

    for i, plan in enumerate(plans):
        x = Inches(0.6 + i * 4.2)
        border = GOLD if plan["recommended"] else None
        bg = LIGHT_BLUE if plan["recommended"] else CARD_BG
        add_rounded_rect(slide, x, Inches(1.7), Inches(3.8), Inches(4.8),
                         fill_color=bg, line_color=border)

        if plan["recommended"]:
            rec_badge = add_rounded_rect(slide, x + Inches(1.2), Inches(1.5),
                                         Inches(1.4), Inches(0.35), fill_color=GOLD)
            add_textbox(slide, x + Inches(1.2), Inches(1.5), Inches(1.4), Inches(0.35),
                        "おすすめ", font_size=12, bold=True, color=WHITE,
                        alignment=PP_ALIGN.CENTER)

        add_textbox(slide, x, Inches(1.9), Inches(3.8), Inches(0.5),
                    plan["name"], font_size=24, bold=True, color=DARK_TEXT,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(2.3), Inches(3.8), Inches(0.3),
                    plan["subtitle"], font_size=14, color=ACCENT_BLUE,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(2.6), Inches(3.8), Inches(0.3),
                    "月額（税別）/ 1〜5名", font_size=11, color=MEDIUM_GRAY,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(2.9), Inches(3.8), Inches(0.7),
                    plan["price"], font_size=36, bold=True, color=GOLD,
                    alignment=PP_ALIGN.CENTER)

        for j, (icon, feat) in enumerate(plan["features"]):
            y = Inches(3.7 + j * 0.38)
            c = DARK_TEXT if icon == "✅" else MEDIUM_GRAY
            add_textbox(slide, x + Inches(0.3), y, Inches(3.2), Inches(0.35),
                        f"{icon}  {feat}", font_size=13, color=c)

        add_rounded_rect(slide, x + Inches(0.3), Inches(5.6), Inches(3.2), Inches(0.7),
                         fill_color=ACCENT_BLUE if plan["recommended"] else LIGHT_GRAY)
        add_textbox(slide, x + Inches(0.3), Inches(5.6), Inches(3.2), Inches(0.7),
                    plan["tagline"], font_size=12, bold=True,
                    color=WHITE if plan["recommended"] else DARK_TEXT,
                    alignment=PP_ALIGN.CENTER)

    # 追加人数の説明
    add_rounded_rect(slide, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.7),
                     fill_color=LIGHT_BLUE, line_color=ACCENT_BLUE)
    add_textbox(slide, Inches(0.8), Inches(6.65), Inches(11.7), Inches(0.6),
                "6名以上の場合：1名追加につき +10,000円/月（税別）  |  すべてのプランは月払い・いつでも変更可能",
                font_size=15, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)


def slide_flow(prs, version):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                "導入の流れ", font_size=36, bold=True, color=DARK_TEXT)
    add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(3))

    steps = [
        ("01", "無料相談", "ボトルネック業務を\nヒアリング"),
        ("02", "ご提案", "御社専用のAI\nエージェント設計書"),
        ("03", "構築", "AI環境の導入\nスキル実装"),
        ("04", "体感", "AIエージェントを\n実際に使ってみる"),
        ("05", "伴走", "月額プランで\n継続的に進化"),
    ]

    for i, (num, title, desc) in enumerate(steps):
        x = Inches(0.3 + i * 2.6)
        # Circle number
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.7),
                                        Inches(1.8), Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT_BLUE
        circle.line.fill.background()
        circle.text_frame.paragraphs[0].text = num
        circle.text_frame.paragraphs[0].font.size = Pt(18)
        circle.text_frame.paragraphs[0].font.bold = True
        circle.text_frame.paragraphs[0].font.color.rgb = WHITE
        circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        if i < 4:
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                         x + Inches(1.6), Inches(1.95),
                                         Inches(0.8), Inches(0.35))
            arr.fill.solid()
            arr.fill.fore_color.rgb = LIGHT_GRAY
            arr.line.fill.background()

        add_textbox(slide, x, Inches(2.7), Inches(2.2), Inches(0.5),
                    title, font_size=18, bold=True, color=GOLD,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(3.2), Inches(2.2), Inches(1.0),
                    desc, font_size=13, color=MEDIUM_GRAY,
                    alignment=PP_ALIGN.CENTER)

    # Phase summary boxes
    if version == "solo":
        build_price = "初期構築プラン：300,000円（期間限定・通常 500,000円）"
    else:
        build_price = "初期構築プラン：150,000円（パートナー特別価格）"

    add_rounded_rect(slide, Inches(0.5), Inches(4.5), Inches(6.0), Inches(2.5),
                     fill_color=CARD_BG)
    add_textbox(slide, Inches(0.7), Inches(4.6), Inches(5.6), Inches(0.4),
                "STEP 1：構築フェーズ", font_size=16, bold=True, color=ACCENT_BLUE)

    step1_items = [
        f"● {build_price}",
        "● AIエージェントがどのようなものかをまず体感",
        "● PC環境のセットアップから外部ツール連携まで完全代行",
    ]
    add_multi_text(slide, Inches(0.7), Inches(5.1), Inches(5.6), Inches(1.8),
                   step1_items, font_size=13, color=MEDIUM_GRAY)

    add_rounded_rect(slide, Inches(6.8), Inches(4.5), Inches(6.0), Inches(2.5),
                     fill_color=LIGHT_BLUE, line_color=ACCENT_BLUE)
    add_textbox(slide, Inches(7.0), Inches(4.6), Inches(5.6), Inches(0.4),
                "STEP 2：伴走フェーズ", font_size=16, bold=True, color=ACCENT_BLUE)

    if version == "solo":
        step2_items = [
            "● 月額 9,800円〜",
            "● 本格的な業務自動化へ",
            "● 最新AI情報の継続提供",
            "● いつでもプラン変更可能",
        ]
    else:
        step2_items = [
            "● 月額 29,800円〜（1〜5名）",
            "● 6名以上：+10,000円/人",
            "● 本格的な業務自動化へ",
            "● いつでもプラン変更可能",
        ]
    add_multi_text(slide, Inches(7.0), Inches(5.1), Inches(5.6), Inches(1.8),
                   step2_items, font_size=13, color=DARK_TEXT)


def slide_cta(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_accent_line(slide, Inches(3.5), Inches(2.0), Inches(6.3))

    add_textbox(slide, Inches(1), Inches(2.3), Inches(11), Inches(1.5),
                "まずは一番手放したい業務を、\n教えてください。",
                font_size=42, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.6),
                "AI SHINEがどのように解決するか、その場でお見せします。",
                font_size=20, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

    add_accent_line(slide, Inches(3.5), Inches(4.8), Inches(6.3))

    add_textbox(slide, Inches(1), Inches(5.3), Inches(11), Inches(0.6),
                "AI SHINEが、あなたの会社の未来を輝かせます。",
                font_size=22, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)

    add_rounded_rect(slide, Inches(3), Inches(6.0), Inches(7.3), Inches(0.7),
                     fill_color=CARD_BG, line_color=GOLD)
    add_textbox(slide, Inches(3), Inches(6.05), Inches(7.3), Inches(0.6),
                "with-AI 株式会社  ｜  contact@with-ai.jp  ｜  090-1296-4814",
                font_size=15, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)


# ====================================================================
# メイン：2つのPPTXを生成
# ====================================================================

def create_deck(version, output_path):
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 共通スライド
    slide_cover(prs, version)           # 1: 表紙
    slide_3walls(prs)                   # 2: 3つの壁
    slide_ai_evolution(prs)             # 3: AI進化
    slide_5year_comparison(prs)         # 4: 5年後比較
    slide_what_is(prs)                  # 5: AI SHINEとは
    slide_3changes(prs)                 # 6: 3つの変化
    slide_architecture(prs)             # 7: 全体像
    slide_withai_orgchart(prs)          # 8: with-AI社の実例組織図
    # 活用事例
    slide_case(prs, 1,
               'スマホでパシャッ！で終わる「経理・バックオフィス自動化」',
               "月末になると領収書整理や請求書発行に追われ、本業の時間が削られる。",
               flow_items=["スマホで\n写真撮影", "Google Drive\nに保存", "AIが自動で\n文字読取", "freeeに\n自動仕訳", "完了！\n確認するだけ"],
               highlight_text="📋 請求書のワンコマンド発行\n\n取引先の名前と金額をチャットで伝えるだけで、登録番号・振込先入りの完璧な請求書 PDFを数秒で自動生成。\n月末は「確認して送信ボタンを押すだけ」に。")
    slide_case(prs, 2,
               '5人の専門家が裏で動く「超速・提案書作成エージェントチーム」',
               "商談後に顧客の業界リサーチや競合調査をしてから提案書を作るため、提出まで何日もかかる。",
               flow_items=["A: 書記\n録音→議事録\n& To-Do抽出", "B: 市場調査\n業界動向・\nペインを調査", "C: 競合調査\n競合サービスを\n並列で分析", "D: 戦略家\nA〜Cを統合し\n提案骨子を設計", "E: 監査役\n批判的検証で\nブラッシュアップ"],
               highlight_text="✅ 結果：あなたが他のお客様対応をしている間に、「お礼メール下書き」と「PowerPoint提案書」が完成して待っています。")
    slide_case(prs, 3,
               'ライバルを丸裸にする「競合サイト・SEO監査エージェント」',
               "新規事業を始めたいが、競合分析やマーケティング調査に回すリソースがない。",
               before_after=(
                   ["▪ 外注に依頼", "▪ 2週間〜1ヶ月待ち", "▪ 費用：30万円〜", "▪ 追加調査にさらに時間と費用"],
                   ["◆ URLを入力するだけ", "◆ わずか数分で完了", "◆ 100点満点のスコア評価", "◆ 改善アクションプラン付き", "◆ プロ級 PDFレポート自動生成"],
               ))
    slide_case(prs, 4,
               'あなたの思考を勝手に発信する「専属ゴーストライター」',
               "集客のためにブログやXを更新したいが、文章を書くのが苦手で続かない。",
               flow_items=["音声メモ・\n気づきを記録", "ネタ帳\nフォルダへ", "AIが文体を\n学習・生成", "図解を\n自動挿入", "WordPress・\nSNSへ下書き"],
               highlight_text="◆ あなたの過去の文章の「癖」や「トーン＆マナー」を学習し、そっくりの文体で自動作成\n◆ 内容に合わせてAIが図解の必要箇所を判断し、画像生成 AIで図解まで自動挿入\n◆ ブログ・X（旧Twitter）のツリー投稿まで、下書きトレイに直接保存")

    slide_skill_system(prs)             # 12: スキルシステム
    slide_comparison_table(prs)         # 13: 比較表
    slide_testimonials(prs)             # 14: 導入企業の声
    slide_step1_build(prs, version)     # 15: STEP1 構築
    slide_campaign(prs, version)        # 16: キャンペーン/価格

    slide_step2_intro(prs)              # 17: STEP2 伴走導入

    # 月額プラン（バージョンで分岐）
    if version == "solo":
        slide_monthly_plans_solo(prs)   # 18: 月額（個人事業主）
    else:
        slide_monthly_plans_corporate(prs)  # 18: 月額（企業向け）

    slide_flow(prs, version)            # 19: 導入の流れ
    slide_cta(prs)                      # 20: CTA

    prs.save(output_path)
    print(f"✅ 生成完了: {output_path}")


if __name__ == "__main__":
    create_deck("solo", "/Users/kaitomain/Desktop/AI_SHINE_個人事業主向け.pptx")
    create_deck("corporate", "/Users/kaitomain/Desktop/AI_SHINE_企業向け.pptx")
    print("\n🎉 2つの資料を生成しました！")
