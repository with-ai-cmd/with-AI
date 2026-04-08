#!/usr/bin/env python3
"""with-AI株式会社 会社紹介PPTX生成スクリプト"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import os

# === 定数 ===
MAIN_DARK = RGBColor(0x0A, 0x3B, 0x8E)
MAIN_LIGHT = RGBColor(0x48, 0xA8, 0xE1)
MAIN_MID = RGBColor(0x29, 0x72, 0xB8)
TEXT_COLOR = RGBColor(0x33, 0x33, 0x33)
SUB_TEXT = RGBColor(0x66, 0x66, 0x66)
ACCENT_BG = RGBColor(0xE8, 0xF4, 0xFD)
LIGHT_GRAY_BG = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_ORANGE = RGBColor(0xF5, 0xA6, 0x23)
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)
FONT_NAME = 'Noto Sans JP'

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
BAR_H = Inches(0.18)

BASE_DIR = os.path.expanduser('~/Desktop/クロードコード')
LOGO_FULL = os.path.join(BASE_DIR, 'img', 'logo-full.png')
LOGO_A = os.path.join(BASE_DIR, 'img', 'logo-A.png')
OUTPUT_PATH = os.path.join(BASE_DIR, 'with-AI_会社紹介.pptx')


# === ヘルパー ===

def set_font(run, size=16, bold=False, color=TEXT_COLOR, name=FONT_NAME):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_text_box(slide, left, top, width, height, text, size=16, bold=False,
                 color=TEXT_COLOR, alignment=PP_ALIGN.LEFT, vertical=MSO_ANCHOR.MIDDLE):
    """テキストボックスを追加（デフォルト上下中央配置）"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = vertical
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color)
    return tf


def add_multiline_text_box(slide, left, top, width, height, lines, alignment=PP_ALIGN.LEFT, vertical=MSO_ANCHOR.MIDDLE):
    """複数行テキスト（各行に異なるスタイル）を追加。linesは(text, size, bold, color)のリスト"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = vertical
    for i, (text, size, bold, color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        run = p.add_run()
        run.text = text
        set_font(run, size=size, bold=bold, color=color)
    return tf


def add_gradient_bar_solid(slide, position='top'):
    """上部/下部にグラデーション風バー（1本の矩形、python-pptx gradient API使用）"""
    if position == 'top':
        top_pos = Emu(0)
    else:
        top_pos = Emu(SLIDE_H) - BAR_H
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), top_pos, SLIDE_W, BAR_H)
    shape.line.fill.background()

    # python-pptx gradient fill API
    fill = shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = RGBColor(0x0A, 0x3B, 0x8E)
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = RGBColor(0x48, 0xA8, 0xE1)
    fill.gradient_stops[1].position = 1.0
    # 方向を左→右(0度)に設定
    fill.gradient_angle = 0


def add_footer_logo(slide):
    if os.path.exists(LOGO_A):
        slide.shapes.add_picture(LOGO_A, Inches(11.5), Inches(6.35), height=Inches(0.7))


def add_rounded_rect(slide, left, top, width, height, fill_color=ACCENT_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.adjustments[0] = 0.04
    return shape


def add_section_title_bar(slide, text):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.55), SLIDE_W, Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BG
    bar.line.fill.background()

    # 左端にグラデーション風アクセント
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.55), Inches(0.12), Inches(1.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = MAIN_DARK
    accent.line.fill.background()

    add_text_box(slide, Inches(0.5), Inches(0.55), Inches(11), Inches(1.15),
                 text, size=32, bold=True, color=MAIN_DARK, alignment=PP_ALIGN.LEFT)


def add_arrow_right(slide, left, top, width=Inches(0.6), height=Inches(0.5), color=MAIN_LIGHT):
    """右向き矢印を追加"""
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    return arrow


def add_chevron(slide, left, top, width=Inches(0.5), height=Inches(0.5), color=MAIN_LIGHT):
    """シェブロン（V字矢印）を追加"""
    chevron = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    chevron.fill.solid()
    chevron.fill.fore_color.rgb = color
    chevron.line.fill.background()
    return chevron


def add_circle_icon(slide, left, top, size_inches, color, text='', text_size=18):
    """円形アイコンを追加"""
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(size_inches), Inches(size_inches))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    if text:
        tf = circle.text_frame
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        set_font(run, size=text_size, bold=True, color=WHITE)
    return circle


def add_decorative_dots(slide, left, top, count=3, spacing=0.18, size=0.08, color=MAIN_LIGHT):
    """装飾的なドットを追加"""
    for i in range(count):
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, left + Inches(i * spacing), top,
            Inches(size), Inches(size)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()


# === スライド作成 ===

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ──────────────────────────────────────
# スライド1: 表紙
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_gradient_bar_solid(slide, 'top')
add_gradient_bar_solid(slide, 'bottom')

# 装飾ドット（右上）
add_decorative_dots(slide, Inches(11.5), Inches(0.7), count=4, spacing=0.22, size=0.1, color=MAIN_LIGHT)
add_decorative_dots(slide, Inches(11.6), Inches(1.0), count=3, spacing=0.22, size=0.1, color=ACCENT_BG)

# ロゴ（中央上部）
if os.path.exists(LOGO_FULL):
    slide.shapes.add_picture(LOGO_FULL, Inches(3.8), Inches(1.2), width=Inches(5.8))

# 装飾ライン
line_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(5.0), Inches(3.5), Inches(3.3), Inches(0.04)
)
line_shape.fill.solid()
line_shape.fill.fore_color.rgb = MAIN_LIGHT
line_shape.line.fill.background()

# タイトル
add_text_box(slide, Inches(2), Inches(3.8), Inches(9.3), Inches(1.2),
             '会社紹介資料', size=44, bold=True, color=TEXT_COLOR, alignment=PP_ALIGN.CENTER)

# サブタイトル
add_text_box(slide, Inches(2), Inches(5.0), Inches(9.3), Inches(0.7),
             'Company Profile', size=22, color=SUB_TEXT, alignment=PP_ALIGN.CENTER)

# 日付
add_text_box(slide, Inches(2), Inches(5.9), Inches(9.3), Inches(0.6),
             '2026年3月', size=16, color=SUB_TEXT, alignment=PP_ALIGN.CENTER)

# 装飾ドット（左下）
add_decorative_dots(slide, Inches(0.8), Inches(6.3), count=3, spacing=0.22, size=0.1, color=MAIN_LIGHT)


# ──────────────────────────────────────
# スライド2: 企業理念
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bar_solid(slide, 'top')
add_gradient_bar_solid(slide, 'bottom')
add_footer_logo(slide)

add_section_title_bar(slide, '企業理念 / Philosophy')

# 装飾（引用マーク風の大きなシェイプ）
quote_mark = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.0), Inches(2.1), Inches(1.3), Inches(0.08))
quote_mark.fill.solid()
quote_mark.fill.fore_color.rgb = MAIN_LIGHT
quote_mark.line.fill.background()

# メインコピー
add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.0),
             '共に次の時代へ。', size=44, bold=True, color=MAIN_DARK, alignment=PP_ALIGN.CENTER)

# 理念テキスト
philosophy_text = (
    '私たちは、クライアントや関わるすべての人を、新しい時代に取り残しません。\n'
    'AIやDXは目的ではなく、人と人とのつながりを守り、\n'
    '人とAIが共存できる環境をつくるための手段です。\n\n'
    'テクノロジーと仕組みによって時間を生み出し、\n'
    '人が本来向き合うべき判断、対話、創造に集中できる状態をつくる。\n'
    'それが、with-AIの変わらない考え方です。'
)
add_text_box(slide, Inches(1.5), Inches(3.7), Inches(10.3), Inches(3.0),
             philosophy_text, size=18, color=TEXT_COLOR, alignment=PP_ALIGN.CENTER,
             vertical=MSO_ANCHOR.TOP)

# 装飾ドット
add_decorative_dots(slide, Inches(0.6), Inches(2.2), count=3, spacing=0.2, size=0.09, color=MAIN_LIGHT)
add_decorative_dots(slide, Inches(12.0), Inches(6.0), count=3, spacing=0.2, size=0.09, color=MAIN_LIGHT)


# ──────────────────────────────────────
# スライド3: ミッション・ビジョン
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bar_solid(slide, 'top')
add_gradient_bar_solid(slide, 'bottom')
add_footer_logo(slide)

add_section_title_bar(slide, 'Mission & Vision')

# ミッションカード
card1 = add_rounded_rect(slide, Inches(0.6), Inches(2.3), Inches(5.8), Inches(4.5), ACCENT_BG)

# ミッションアイコン
add_circle_icon(slide, Inches(1.2), Inches(2.6), 0.6, MAIN_DARK, '▶', 16)

add_text_box(slide, Inches(2.0), Inches(2.6), Inches(4.0), Inches(0.6),
             'Mission', size=20, bold=True, color=MAIN_LIGHT)
add_text_box(slide, Inches(1.2), Inches(3.5), Inches(4.8), Inches(1.0),
             'ともに考え、ともに整え、\n次の時代へ進む。', size=26, bold=True, color=MAIN_DARK)
add_text_box(slide, Inches(1.2), Inches(4.7), Inches(4.8), Inches(1.8),
             'AIや手法を目的にするのではなく、\n人と現場が主役のまま、\n変化に向き合える状態をつくる。\nそれが、with-AIの使命です。',
             size=16, color=TEXT_COLOR, vertical=MSO_ANCHOR.TOP)

# 中央の矢印
add_arrow_right(slide, Inches(6.55), Inches(4.3), Inches(0.6), Inches(0.5), MAIN_LIGHT)

# ビジョンカード
card2 = add_rounded_rect(slide, Inches(6.9), Inches(2.3), Inches(5.8), Inches(4.5), ACCENT_BG)

# ビジョンアイコン
add_circle_icon(slide, Inches(7.5), Inches(2.6), 0.6, MAIN_LIGHT, '◎', 16)

add_text_box(slide, Inches(8.3), Inches(2.6), Inches(4.0), Inches(0.6),
             'Vision', size=20, bold=True, color=MAIN_LIGHT)
add_text_box(slide, Inches(7.5), Inches(3.5), Inches(4.8), Inches(1.0),
             '人が判断し、\n仕組みが支える時代へ。', size=26, bold=True, color=MAIN_DARK)
add_text_box(slide, Inches(7.5), Inches(4.7), Inches(4.8), Inches(1.8),
             'テクノロジーが進化しても、\n決めるのは人であり、支えるのは仕組み。\nそのバランスが取れた未来を、\n私たちは実装します。',
             size=16, color=TEXT_COLOR, vertical=MSO_ANCHOR.TOP)


# ──────────────────────────────────────
# スライド4: バリュー
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bar_solid(slide, 'top')
add_gradient_bar_solid(slide, 'bottom')
add_footer_logo(slide)

add_section_title_bar(slide, 'Values — 私たちが大切にしていること')

values = [
    ('01', '人を主役にする', 'AIを入れることが目的ではなく、\n人がより良く働き、より良く\n判断できる状態をつくることを\n優先する。'),
    ('02', '現場から考える', '机上の理想ではなく、\n現場で実際に機能する形まで\n落とし込む。'),
    ('03', '取り残さず、\nともに進む', '変化に不安を感じる人がいても\n置いていかない。\n理解し、寄り添い、伴走しながら\n次の時代へ進む。'),
    ('04', '仕組みで支える', '属人化や非効率を放置せず、\n再現性のある仕組みに\n変えていく。'),
    ('05', 'つながりを\n大切にする', 'デジタル時代だからこそ、\n人と人との関係性を\n大切にする。'),
]

for i, (num, title, desc) in enumerate(values):
    left = Inches(0.3 + i * 2.58)
    top = Inches(2.3)

    # カード背景
    card = add_rounded_rect(slide, left, top, Inches(2.35), Inches(4.5), ACCENT_BG)

    # 番号の円（グラデーション色）
    ratio = i / 4
    r = int(0x0A + (0x48 - 0x0A) * ratio)
    g = int(0x3B + (0xA8 - 0x3B) * ratio)
    b = int(0x8E + (0xE1 - 0x8E) * ratio)
    add_circle_icon(slide, left + Inches(0.78), top + Inches(0.3), 0.8, RGBColor(r, g, b), num, 20)

    # タイトル
    add_text_box(slide, left + Inches(0.1), top + Inches(1.3), Inches(2.15), Inches(0.8),
                 title, size=17, bold=True, color=MAIN_DARK, alignment=PP_ALIGN.CENTER)

    # 区切りライン
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  left + Inches(0.6), top + Inches(2.2),
                                  Inches(1.15), Inches(0.03))
    sep.fill.solid()
    sep.fill.fore_color.rgb = MAIN_LIGHT
    sep.line.fill.background()

    # 説明
    add_text_box(slide, left + Inches(0.15), top + Inches(2.5), Inches(2.05), Inches(1.8),
                 desc, size=13, color=TEXT_COLOR, alignment=PP_ALIGN.CENTER,
                 vertical=MSO_ANCHOR.TOP)


# ──────────────────────────────────────
# スライド5: 会社概要
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bar_solid(slide, 'top')
add_gradient_bar_solid(slide, 'bottom')
add_footer_logo(slide)

add_section_title_bar(slide, '会社概要 / Company Overview')

# ロゴを左上に配置
if os.path.exists(LOGO_FULL):
    slide.shapes.add_picture(LOGO_FULL, Inches(9.5), Inches(0.65), height=Inches(0.8))

info_items = [
    ('会社名', 'with-AI株式会社'),
    ('英語表記', 'with-AI Inc.'),
    ('設立', '2025年10月'),
    ('所在地', '東京都渋谷区千駄ヶ谷5-16-10'),
    ('代表者', '勝又海斗（代表取締役）'),
    ('従業員数', '10名'),
    ('資本金', '300万円'),
    ('事業内容', '社外AI顧問 / AIエージェント作成 / AI教育・研修 / with-AI Plus / AI採用支援 / コミュニティ連携 / システム開発'),
]

row_h = Inches(0.58)
for i, (label, value) in enumerate(info_items):
    y = Inches(2.2 + i * 0.62)

    # 交互の背景色
    if i % 2 == 0:
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), y, Inches(10.9), row_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = LIGHT_GRAY_BG
        bg.line.fill.background()

    # 左端にアクセントライン
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), y, Inches(0.06), row_h)
    accent.fill.solid()
    accent.fill.fore_color.rgb = MAIN_DARK
    accent.line.fill.background()

    # ラベル
    add_text_box(slide, Inches(1.5), y, Inches(2.2), row_h,
                 label, size=16, bold=True, color=MAIN_DARK)

    # 値
    add_text_box(slide, Inches(4.0), y, Inches(7.8), row_h,
                 value, size=16, color=TEXT_COLOR)

# 装飾ドット（右下）
add_decorative_dots(slide, Inches(0.6), Inches(6.8), count=4, spacing=0.2, size=0.09, color=MAIN_LIGHT)


# ──────────────────────────────────────
# スライド6: 事業紹介（概要）
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bar_solid(slide, 'top')
add_gradient_bar_solid(slide, 'bottom')
add_footer_logo(slide)

add_section_title_bar(slide, '事業紹介 / Our Services')

services = [
    ('社外AI顧問', '経営者に寄り添い、AI活用の\n戦略立案から仕組み化まで伴走支援', '🤝'),
    ('AIKOMON', '経営者の知識・判断基準を\nAIエージェントとして再現', '🧠'),
    ('AI教育・研修', '実務起点の研修で\n現場定着まで設計', '📚'),
    ('with-AI Plus', '実務直結型のAI学習\nプラットフォーム', '💻'),
    ('AI採用支援', '相互理解型の採用で\n人材の魅力を可視化', '👥'),
    ('コミュニティ連携', '経営者ネットワークを通じた\n案件創出と価値提供', '🌐'),
]

for i, (title, desc, icon) in enumerate(services):
    row = i // 3
    col = i % 3
    left = Inches(0.5 + col * 4.2)
    top = Inches(2.3 + row * 2.7)

    card = add_rounded_rect(slide, left, top, Inches(3.85), Inches(2.35), ACCENT_BG)

    # 番号バッジ（グラデーション色）
    ratio = i / 5
    r = int(0x0A + (0x48 - 0x0A) * ratio)
    g = int(0x3B + (0xA8 - 0x3B) * ratio)
    b = int(0x8E + (0xE1 - 0x8E) * ratio)
    badge_color = RGBColor(r, g, b)

    add_circle_icon(slide, left + Inches(0.25), top + Inches(0.25), 0.55, badge_color, f'0{i+1}', 14)

    # タイトル
    add_text_box(slide, left + Inches(0.95), top + Inches(0.25), Inches(2.6), Inches(0.55),
                 title, size=20, bold=True, color=MAIN_DARK)

    # 区切りライン
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  left + Inches(0.25), top + Inches(1.0),
                                  Inches(3.35), Inches(0.025))
    sep.fill.solid()
    sep.fill.fore_color.rgb = MAIN_LIGHT
    sep.line.fill.background()

    # 説明
    add_text_box(slide, left + Inches(0.25), top + Inches(1.15), Inches(3.35), Inches(1.1),
                 desc, size=15, color=TEXT_COLOR, vertical=MSO_ANCHOR.TOP)

# 行間の矢印（1行目→2行目）
for col in range(3):
    x = Inches(2.2 + col * 4.2)
    slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x, Inches(4.72), Inches(0.3), Inches(0.35)).fill.solid()
    arrow = slide.shapes._spTree[-1]
    slide.shapes[-1].fill.fore_color.rgb = MAIN_LIGHT
    slide.shapes[-1].line.fill.background()


# ──────────────────────────────────────
# スライド7: AIKOMONの詳細
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bar_solid(slide, 'top')
add_gradient_bar_solid(slide, 'bottom')
add_footer_logo(slide)

add_section_title_bar(slide, 'AIKOMON — 経営者専用AIエージェント')

# サービス概要
add_text_box(slide, Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.2),
             '社長個人の頭の中にある知識、経験、判断基準、業務の進め方を言語化・整理し、\nAIエージェントとして再現。経営者の専属サポートチームのように機能するAI環境を構築します。',
             size=18, color=TEXT_COLOR, vertical=MSO_ANCHOR.TOP)

# フロー矢印（プロセス）
flow_items = ['ヒアリング', '→', '知識の整理・構造化', '→', 'AIエージェント構築', '→', '運用・改善']
flow_left = Inches(0.8)
for j, item in enumerate(flow_items):
    w = Inches(2.2) if j % 2 == 0 else Inches(0.5)
    if j % 2 == 0:
        box = add_rounded_rect(slide, flow_left, Inches(3.5), w, Inches(0.5), MAIN_DARK)
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = item
        set_font(run, size=13, bold=True, color=WHITE)
    else:
        add_text_box(slide, flow_left, Inches(3.45), w, Inches(0.55),
                     '▶', size=20, bold=True, color=MAIN_LIGHT, alignment=PP_ALIGN.CENTER)
    flow_left += w

# 導入効果カード
effect_icons = ['⏱', '🔄', '🚀']
effects = [
    ('スキマ時間の有効活用', '移動中・商談の合間・出先など、\n場所を問わずAIエージェントを\n活用可能'),
    ('思考・判断の再現性向上', '社長の感覚や経験に依存していた\n判断を整理し、基準を明確化。\n指示の再現性が高まる'),
    ('経営者の時間創出', 'AIが情報整理・たたき台作成・\n論点整理を担い、社長は経営判断・\n対話・戦略設計に集中可能'),
]

for i, (title, desc) in enumerate(effects):
    left = Inches(0.5 + i * 4.2)
    top = Inches(4.3)

    card = add_rounded_rect(slide, left, top, Inches(3.85), Inches(2.7), ACCENT_BG)

    # アイコン風の円
    ratio = i / 2
    r = int(0x0A + (0x48 - 0x0A) * ratio)
    g = int(0x3B + (0xA8 - 0x3B) * ratio)
    b = int(0x8E + (0xE1 - 0x8E) * ratio)
    add_circle_icon(slide, left + Inches(0.25), top + Inches(0.25), 0.5, RGBColor(r, g, b), f'0{i+1}', 14)

    # タイトル
    add_text_box(slide, left + Inches(0.9), top + Inches(0.25), Inches(2.7), Inches(0.5),
                 title, size=18, bold=True, color=MAIN_DARK)

    # 区切りライン
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  left + Inches(0.25), top + Inches(0.95),
                                  Inches(1.8), Inches(0.03))
    sep.fill.solid()
    sep.fill.fore_color.rgb = MAIN_LIGHT
    sep.line.fill.background()

    # 説明
    add_text_box(slide, left + Inches(0.25), top + Inches(1.2), Inches(3.35), Inches(1.3),
                 desc, size=15, color=TEXT_COLOR, vertical=MSO_ANCHOR.TOP)


# ──────────────────────────────────────
# スライド8: 強み・差別化
# ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bar_solid(slide, 'top')
add_gradient_bar_solid(slide, 'bottom')
add_footer_logo(slide)

add_section_title_bar(slide, 'with-AIの強み / Our Strengths')

# ── POINT 01 ──
card1 = add_rounded_rect(slide, Inches(0.5), Inches(2.3), Inches(6.0), Inches(4.7), ACCENT_BG)

# POINT番号バッジ
add_circle_icon(slide, Inches(0.9), Inches(2.6), 0.65, MAIN_DARK, '01', 18)

add_text_box(slide, Inches(1.7), Inches(2.6), Inches(4.5), Inches(0.65),
             'POINT', size=16, bold=True, color=MAIN_LIGHT)
add_text_box(slide, Inches(0.9), Inches(3.4), Inches(5.2), Inches(1.0),
             '会社ごとの\n"使える仕組み"までつくる', size=24, bold=True, color=MAIN_DARK)

# 比較表示（矢印付き）
# 他社
other_box = add_rounded_rect(slide, Inches(1.0), Inches(4.5), Inches(5.0), Inches(0.7), LIGHT_GRAY_BG)
add_text_box(slide, Inches(1.2), Inches(4.5), Inches(4.8), Inches(0.7),
             '他社 ▶ 「このツールが便利です」で終わることが多い',
             size=14, color=SUB_TEXT)

# 矢印（下向き）
down_arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(3.3), Inches(5.3), Inches(0.35), Inches(0.35))
down_arrow.fill.solid()
down_arrow.fill.fore_color.rgb = ACCENT_ORANGE
down_arrow.line.fill.background()

# with-AI
our_box = add_rounded_rect(slide, Inches(1.0), Inches(5.75), Inches(5.0), Inches(0.95), WHITE)
# 左端アクセント
left_acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(5.75), Inches(0.08), Inches(0.95))
left_acc.fill.solid()
left_acc.fill.fore_color.rgb = MAIN_DARK
left_acc.line.fill.background()

add_text_box(slide, Inches(1.3), Inches(5.75), Inches(4.6), Inches(0.95),
             'with-AI ▶ 会社ごとの業務・社長の考え方・現場の\n流れに合わせて、業務設計と仕組み化まで行う',
             size=14, bold=True, color=MAIN_DARK)

# ── POINT 02 ──
card2 = add_rounded_rect(slide, Inches(6.85), Inches(2.3), Inches(6.0), Inches(4.7), ACCENT_BG)

add_circle_icon(slide, Inches(7.25), Inches(2.6), 0.65, MAIN_LIGHT, '02', 18)

add_text_box(slide, Inches(8.05), Inches(2.6), Inches(4.5), Inches(0.65),
             'POINT', size=16, bold=True, color=MAIN_LIGHT)
add_text_box(slide, Inches(7.25), Inches(3.4), Inches(5.2), Inches(1.0),
             '研修会社でも開発会社でもなく\nその間を埋める存在', size=24, bold=True, color=MAIN_DARK)

# 研修会社
r_box = add_rounded_rect(slide, Inches(7.35), Inches(4.5), Inches(2.3), Inches(0.6), LIGHT_GRAY_BG)
add_text_box(slide, Inches(7.45), Inches(4.5), Inches(2.1), Inches(0.6),
             '研修会社\n教えるのが得意', size=12, color=SUB_TEXT, alignment=PP_ALIGN.CENTER)

# 矢印（右向き→統合）
add_arrow_right(slide, Inches(9.75), Inches(4.55), Inches(0.4), Inches(0.45), ACCENT_ORANGE)

# 開発会社
d_box = add_rounded_rect(slide, Inches(10.25), Inches(4.5), Inches(2.3), Inches(0.6), LIGHT_GRAY_BG)
add_text_box(slide, Inches(10.35), Inches(4.5), Inches(2.1), Inches(0.6),
             '開発会社\n作るのが得意', size=12, color=SUB_TEXT, alignment=PP_ALIGN.CENTER)

# 統合矢印（下向き）
down_arrow2 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(9.7), Inches(5.2), Inches(0.35), Inches(0.4))
down_arrow2.fill.solid()
down_arrow2.fill.fore_color.rgb = ACCENT_ORANGE
down_arrow2.line.fill.background()

# with-AI（その間を埋める）
our_box2 = add_rounded_rect(slide, Inches(7.35), Inches(5.75), Inches(5.2), Inches(0.95), WHITE)
left_acc2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.35), Inches(5.75), Inches(0.08), Inches(0.95))
left_acc2.fill.solid()
left_acc2.fill.fore_color.rgb = MAIN_DARK
left_acc2.line.fill.background()

add_text_box(slide, Inches(7.6), Inches(5.75), Inches(4.8), Inches(0.95),
             'with-AI ▶ 何を整理すべきか、どう業務に落とすか、\n誰がどう使うか、どう定着させるかまでつなげる',
             size=14, bold=True, color=MAIN_DARK)


# ──────────────────────────────────────
# 保存
# ──────────────────────────────────────
prs.save(OUTPUT_PATH)
print(f'✅ 生成完了: {OUTPUT_PATH}')
print(f'   スライド数: {len(prs.slides)}枚')
